# dependencies
# pip install python-pptx
# pip install glob2

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob

path = "c:/Users/lozas/COMM1076-python/"
testString = "Learning Objectives"

print ("---- begin ----")
for eachfile in glob.glob(path + "*.pptx"):
    prs = Presentation(eachfile)
    #print(eachfile)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find("Chapter"))!=-1:
                    print(shape.text)
                if(shape.text.find(testString))!=-1:
                    #print("id: %s, type: %s" % (shape.shape_id, shape.shape_type))
                    print(shape.text)